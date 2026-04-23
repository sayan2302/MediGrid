from pptx import Presentation
from pptx.util import Inches, Pt
import sys

def create_presentation():
    prs = Presentation()

    # Define slide layouts
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # Slide 1: Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "MediGrid: Supply Chain Management in Hospitals"
    subtitle.text = "Sayan Pramanick, BITS Pilani\nDissertation Work at Tipstat Infotech Ltd."

    # Slide 2: Problem Statement
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Problem Statement"
    tf = body_shape.text_frame
    tf.text = "Hospital supply chains face significant operational challenges:"
    p = tf.add_paragraph()
    p.text = "Frequent stock-outs, overstocking, and expiry-related wastage."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Lack of real-time inventory visibility across departments."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Existing systems lack intelligent decision-support and predictive analytics."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Poor coordination and reliance on manual procurement processes."
    p.level = 1

    # Slide 3: The Proposed Solution: MediGrid
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "The Proposed Solution: MediGrid"
    tf = body_shape.text_frame
    tf.text = "A modular, scalable solution for hospital supply chain management:"
    p = tf.add_paragraph()
    p.text = "Real-time, granular inventory tracking at the batch level."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Automated procurement workflows and comprehensive vendor management."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Intelligent AI-driven demand forecasting and expiry risk assessment."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Built to provide analytics-driven reporting and actionable decision support."
    p.level = 1

    # Slide 4: System Architecture
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "System Architecture"
    tf = body_shape.text_frame
    tf.text = "MediGrid utilizes a modern multi-tiered architecture:"
    p = tf.add_paragraph()
    p.text = "Frontend (React/Vite): Dynamic, visually rich KPI dashboards and workflows."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Backend (Node/Express/MongoDB): Core logic handling inventory transactions and APIs."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "AI Service (FastAPI/Python): Python microservice utilizing Llama 3.1 8B via Groq."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Fuses traditional operations research formulas with LLM inference."
    p.level = 1

    # Slide 5: Inventory Optimization & Logic
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Inventory Optimization & Logic"
    tf = body_shape.text_frame
    tf.text = "To prevent stock-outs, MediGrid uses Reorder Point Models:"
    p = tf.add_paragraph()
    p.text = "Formula: RP = (Average Demand × Lead Time) + Safety Stock"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Automated Tracking: The backend monitors inventory against the RP in real-time."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Alert System: Generates WARNING or CRITICAL alerts when stock drops below threshold."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Triggers automated workflows to map vendors and accelerate procurement."
    p.level = 1

    # Slide 6: AI-Driven Demand Forecasting
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "AI-Driven Demand Forecasting"
    tf = body_shape.text_frame
    tf.text = "A dual-layered approach to predicting medical supply demand:"
    p = tf.add_paragraph()
    p.text = "Fallback Heuristic (Moving Average): Forecast = (D1 + D2 + ... + Dn) / n"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "LLM-Driven Inference: Evaluates historical usage sequences."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Output: AI generates predicted demand, confidence score, and analytical reasoning."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Provides robust predictions even with fluctuating consumption patterns."
    p.level = 1

    # Slide 7: Expiry Risk Evaluation
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Expiry Risk Evaluation"
    tf = body_shape.text_frame
    tf.text = "Intelligently preventing wastage through batch-level tracking:"
    p = tf.add_paragraph()
    p.text = "Coverage Calculation: coverage_days = currentStock / dailyUsageRate"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "AI Assessment: Evaluates current stock, daily usage, and days to expiry."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Determines risk level (HIGH, MEDIUM, LOW)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Example: Flags high risk if coverage days significantly exceed days to expiry."
    p.level = 1

    # Slide 8: Demo Walkthrough - Demand Forecasting
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Demo Walkthrough: Demand Forecasting"
    tf = body_shape.text_frame
    tf.text = "Scenario: Predicting demand for Paracetamol 500mg over 14 days."
    p = tf.add_paragraph()
    p.text = "User inputs horizon days in the AI Insights Dashboard."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Backend forwards usage history to Python AI service."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "LLM returns structured JSON with predicted values and reasoning."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "UI dynamically renders trend charts, KPI cards, and the AI Verdict."
    p.level = 1

    # Slide 9: Demo Walkthrough - Expiry & Alerts
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Demo Walkthrough: Expiry & Alerts"
    tf = body_shape.text_frame
    tf.text = "Scenario: Amoxicillin stock and ICU surgical mask alerts."
    p = tf.add_paragraph()
    p.text = "Expiry Risk: Amoxicillin consumption is evaluated against its expiration date."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Visual burndown charts clearly highlight potential wastage (HIGH RISK)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Procurement Alert: Pulling surgical masks drops stock below the Reorder Point."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Triggers an alert, allowing immediate generation of a Purchase Order to mapped vendors."
    p.level = 1

    # Slide 10: Conclusion & Future Work
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Conclusion & Future Work"
    tf = body_shape.text_frame
    tf.text = "MediGrid successfully modernizes hospital logistics:"
    p = tf.add_paragraph()
    p.text = "Provides a scalable foundation for intelligent healthcare supply chains."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Improves KPI metrics: stock-out reduction, turnover ratio, and expiry loss."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Future Work: Deep integration with existing ERP systems."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Future Work: Expanding AI models for advanced stochastic time-series forecasting."
    p.level = 1

    # Save presentation
    prs.save("MediGrid_Presentation.pptx")
    print("Presentation saved as MediGrid_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
